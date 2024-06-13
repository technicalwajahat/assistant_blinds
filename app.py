from io import BytesIO

import docx2txt
import uvicorn
from PyPDF2 import PdfReader
from fastapi import FastAPI, File, UploadFile, HTTPException
from pptx import Presentation
from gtts import gTTS

app = FastAPI()


# Function to convert text to audio using gTTS
def text_to_audio(text, language='en'):
    tts = gTTS(text=text, lang=language, slow=False)
    fp = BytesIO()
    tts.write_to_fp(fp)
    fp.seek(0)
    return fp


# Docx to Audio
@app.post("/convert-docx-to-audio")
async def convert_to_audio(file: UploadFile = File(...)):
    if file.content_type not in {"application/msword"}:
        raise HTTPException(status_code=415, detail="Unsupported file format. Supported formats: DOC/DOCX")

    try:
        file_content = await file.read()
        with BytesIO(file_content) as docx_file:
            text = docx2txt.process(docx_file)

        text = text.replace('\n', " ").replace('\t', "")
        return text

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error converting file: {str(e)}")


# PDF to Audio
@app.post("/convert-pdf-to-audio")
async def convert_pdf_to_audio(file: UploadFile = File(...)):
    if file.content_type not in "application/pdf":
        raise HTTPException(status_code=415, detail="Unsupported file format. Supported format: PDF")

    try:
        file_content = await file.read()
        pdf_reader = PdfReader(BytesIO(file_content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()

        text = text.replace('\n', " ").replace('\t', "")
        return text

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error converting PDF file: {str(e)}")


# PPTX to Audio
@app.post("/convert-pptx-to-audio")
async def convert_pptx_to_audio(file: UploadFile = File(...)):
    if file.content_type not in {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}:
        raise HTTPException(status_code=415, detail="Unsupported file format. Supported format: PPTX")

    try:
        file_content = await file.read()
        with BytesIO(file_content) as pptx_file:
            presentation = Presentation(pptx_file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        text = text.replace('\n', " ").replace('\t', "")
        return text
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error converting PPTX file: {str(e)}")


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
